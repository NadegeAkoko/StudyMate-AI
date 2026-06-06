from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import re


def clean_text(text):

    # supprimer les retours à la ligne multiples
    text = re.sub(r"\n+", "\n", text)

    # supprimer les espaces multiples
    text = re.sub(r"\s+", " ", text)

    # supprimer les références HAL
    text = re.sub(r"hal-\d+", "", text, flags=re.IGNORECASE)

    # supprimer les années isolées
    text = re.sub(r"\b(19|20)\d{2}\b", "", text)

    return text.strip()


def load_pdf(uploaded_file):

    text = ""

    pdf = PdfReader(uploaded_file)

    for page in pdf.pages:

        page_text = page.extract_text()

        if page_text:
            text += page_text + "\n"

    text = clean_text(text)

    return text


def load_docx(uploaded_file):

    doc = Document(uploaded_file)

    text = "\n".join(
        paragraph.text
        for paragraph in doc.paragraphs
    )

    text = clean_text(text)

    return text


def load_pptx(uploaded_file):

    presentation = Presentation(uploaded_file)

    text = ""

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):
                text += shape.text + "\n"

    text = clean_text(text)

    return text